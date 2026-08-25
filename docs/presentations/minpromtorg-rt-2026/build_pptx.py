#!/usr/bin/env python3
"""Сборка PPTX для Минпромторга РТ. Запуск: python3 build_pptx.py"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("KD-Minpromtorg-RT-2026.pptx")
LOGO = Path(__file__).with_name("logo.png")

INK = RGBColor(0x12, 0x24, 0x1B)
GREEN = RGBColor(0x1B, 0x5E, 0x3B)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
CREAM = RGBColor(0xF6, 0xF1, 0xE6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x5A, 0x6B, 0x60)
DEEP = RGBColor(0x0C, 0x24, 0x18)
CARD = RGBColor(0xFF, 0xFD, 0xF8)


def _set_run(run, size, bold=False, color=INK, name="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = name
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        # keep East Asian fallback
        pass


def box(slide, l, t, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.adjustments[0] = 0.08
    sh.line.fill.background()
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.fill.solid()
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    return sh


def text(slide, l, t, w, h, content, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(content, str):
        content = [content]
    for i, line in enumerate(content):
        para = p if i == 0 else tf.add_paragraph()
        para.alignment = align
        run = para.add_run()
        run.text = line
        _set_run(run, size, bold, color)
    return tb


def footer(slide, n, total=13, dark=False):
    color = RGBColor(0x9A, 0xA8, 0x9E) if not dark else RGBColor(0x7A, 0x8C, 0x80)
    text(
        slide, Inches(0.6), Inches(7.15), Inches(10.5), Inches(0.28),
        "ООО «Казанские Деликатесы»  ·  ИНН 1686021074  ·  для служебного пользования",
        size=11, color=color,
    )
    text(
        slide, Inches(11.6), Inches(7.15), Inches(1.1), Inches(0.28),
        f"{n} / {total}", size=11, color=color, align=PP_ALIGN.RIGHT,
    )


def blank(prs, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = DEEP if dark else CREAM
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    # send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return slide


def add_logo(slide, dark=False):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(0.6), Inches(0.28), Inches(0.42), Inches(0.42))
    text(
        slide, Inches(1.12), Inches(0.32), Inches(6), Inches(0.36),
        "ООО «Казанские Деликатесы»",
        size=14, bold=True, color=CREAM if dark else GREEN,
    )


def kpi(slide, l, t, w, h, num, lbl, dark=False):
    box(slide, l, t, w, h, fill=DEEP if dark else WHITE, line=None if dark else RGBColor(0xE4, 0xDD, 0xD0))
    text(slide, l + Inches(0.16), t + Inches(0.14), w - Inches(0.3), Inches(0.46),
         num, size=22, bold=True, color=GOLD if dark else GREEN)
    text(slide, l + Inches(0.16), t + Inches(0.6), w - Inches(0.3), Inches(0.55),
         lbl, size=12, color=CREAM if dark else MUTED)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 title
    s = blank(prs, dark=True)
    add_logo(s, dark=True)
    text(s, Inches(0.6), Inches(1.3), Inches(12), Inches(0.35),
         "ДЛЯ МИНИСТЕРСТВА ПРОМЫШЛЕННОСТИ И ТОРГОВЛИ РЕСПУБЛИКИ ТАТАРСТАН",
         size=13, bold=True, color=GOLD)
    text(s, Inches(0.6), Inches(1.8), Inches(11.5), Inches(2.2),
         ["Производство халяльной", "мясной продукции.", "Роботизация линий."],
         size=40, bold=True, color=CREAM)
    text(s, Inches(0.6), Inches(4.2), Inches(10), Inches(0.6),
         "Казань, Агропромышленный парк. Выступление по проекту модернизации и мерам поддержки 2026 года.",
         size=16, color=CREAM)
    kpi(s, Inches(0.6), Inches(5.15), Inches(2.9), Inches(1.25), "869 млн ₽", "выручка 2025, БФО ФНС", True)
    kpi(s, Inches(3.7), Inches(5.15), Inches(2.9), Inches(1.25), "> 1 млрд ₽", "прогноз оборота 2026", True)
    kpi(s, Inches(6.8), Inches(5.15), Inches(2.9), Inches(1.25), "110", "сотрудников сейчас", True)
    kpi(s, Inches(9.9), Inches(5.15), Inches(2.8), Inches(1.25), "3 600 м²", "производство, склад, холод", True)
    footer(s, 1, dark=True)

    # 2 snapshot
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "СНИМОК ПРЕДПРИЯТИЯ", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.55), "Что мы имеем сегодня", 32, True, INK)
    text(s, Inches(0.6), Inches(1.7), Inches(12), Inches(0.4),
         "Действующий завод полного цикла в Казани. Площади и штат — факт на август 2026.", 16, False, MUTED)
    cards = [
        ("110", "человек в штате. ССЧ ФНС 2025: 64"),
        ("2 000 м²", "производственные помещения"),
        ("1 300 м²", "холодильные камеры"),
        ("300 м²", "складские помещения"),
        ("64 SKU", "пепперони, сосиски, ветчины, казылык, выпечка"),
        ("2022", "регистрация. В 2023 площадку посетил Раис РТ Р. Н. Минниханов"),
        ("ОСНО", "ИНН 1686021074 · ОКВЭД 10.13 · Агропромпарк «Казань»"),
        ("3 600 м²", "суммарно производство + холод + склад"),
    ]
    for i, (n, l) in enumerate(cards):
        col, row = i % 4, i // 4
        kpi(s, Inches(0.6 + col * 3.15), Inches(2.3 + row * 2.15), Inches(3.0), Inches(1.95), n, l)
    footer(s, 2)

    # 3 about
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "О КОМПАНИИ", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.6), "Халяльный производитель Республики Татарстан", 28, True, INK)
    bullets = [
        "ООО «Казанские Деликатесы». Гендиректор — Адаева Гульнара Нурхаматовна.",
        "Сырьё: говядина, конина, курица, индейка. Свинины нет.",
        "Халяль ДУМ РТ № 614A/2024. HACCP. ISO 22000:2018. ТР ТС 021/2011.",
        "Контрактное производство / СТМ: линейка Aslam для АО «ОМПК».",
        "Поставки: Татарстан, федеральные сети, АЗС, HoReCa, экспорт ЕАЭС.",
    ]
    y = Inches(1.95)
    for b in bullets:
        box(s, Inches(0.6), y, Inches(7.6), Inches(0.72), WHITE, RGBColor(0xE4, 0xDD, 0xD0))
        text(s, Inches(0.8), y + Inches(0.16), Inches(7.2), Inches(0.5), b, 15, False, INK)
        y += Inches(0.82)
    box(s, Inches(8.5), Inches(1.95), Inches(4.2), Inches(4.4), WHITE, RGBColor(0xE4, 0xDD, 0xD0))
    text(s, Inches(8.7), Inches(2.15), Inches(3.8), Inches(0.35), "РЕКВИЗИТЫ", 12, True, GREEN)
    text(s, Inches(8.7), Inches(2.6), Inches(3.8), Inches(3.4),
         ["ИНН 1686021074", "ОГРН 1221600096893", "Казань, ул. Аграрная, 2, оф. 7",
          "+7 987 217-02-02", "info@kazandelikates.tatar", "kazandelikates.tatar"],
         16, False, INK)
    footer(s, 3)

    # 4 revenue
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "ФИНАНСЫ · ВЫРУЧКА", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.6), "От 59 млн до 869 млн за три отчётных года", 28, True, INK)
    years = [("2023", "58,6", 0.18), ("2024  +435%", "313,9", 0.40), ("2025  +177%", "869,1", 0.82), ("2026 прогноз", "> 1 000", 0.94)]
    for i, (yr, val, h) in enumerate(years):
        x = Inches(1.2 + i * 3.0)
        height = Inches(3.4 * h)
        top = Inches(5.35) - height
        sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, Inches(2.1), height)
        sh.adjustments[0] = 0.08
        sh.line.fill.background()
        sh.fill.solid()
        sh.fill.fore_color.rgb = GOLD if i == 3 else GREEN
        text(s, x, top - Inches(0.38), Inches(2.1), Inches(0.35), val, 16, True, INK, PP_ALIGN.CENTER)
        text(s, x, Inches(5.45), Inches(2.1), Inches(0.4), yr, 13, True, MUTED, PP_ALIGN.CENTER)
    text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.45),
         "млн ₽, стр. 2110 БФО. 2023–2025 — ФНС / ГИР БО. 2026 — прогноз: превысить 1 млрд ₽. 1 кв. 2026 (управленка): 229 млн ₽.",
         12, False, MUTED)
    footer(s, 4)

    # 5 table-like
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "ФИНАНСЫ · РЕЗУЛЬТАТ", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.55), "Прибыль растёт вместе с масштабом", 28, True, INK)
    rows = [
        ("Показатель", "2023", "2024", "2025"),
        ("Выручка", "58,6 млн", "313,9 млн", "869,1 млн"),
        ("Чистая прибыль", "0,13 млн", "39,6 млн", "68,9 млн"),
        ("Норма чистой прибыли", "0,2%", "12,6%", "7,9%"),
        ("Активы на 31 декабря", "—", "146 млн*", "266 млн"),
        ("Собственный капитал", "0,14 млн", "39,8 млн", "96,2 млн"),
        ("Основные средства", "—", "38,7 млн*", "68,9 млн"),
        ("ССЧ, ФНС", "9", "32", "64"),
    ]
    y = Inches(1.85)
    for i, row in enumerate(rows):
        bg = WHITE if i % 2 == 0 else RGBColor(0xEE, 0xE8, 0xDA)
        box(s, Inches(0.6), y, Inches(12.1), Inches(0.48), bg)
        text(s, Inches(0.75), y + Inches(0.08), Inches(4.4), Inches(0.35), row[0], 14, i == 0, MUTED if i == 0 else INK)
        for j, cell in enumerate(row[1:]):
            text(s, Inches(5.4 + j * 2.4), y + Inches(0.08), Inches(2.2), Inches(0.35),
                 cell, 14, i == 0 or j == 2, GREEN if j == 2 and i else (MUTED if i == 0 else INK), PP_ALIGN.RIGHT)
        y += Inches(0.48)
    text(s, Inches(0.6), Inches(6.55), Inches(12), Inches(0.4),
         "*2024 активы и ОС — оценка по темпам Checko 2025. Рентабельность продаж 2024: 15,9% (выше ¾ отрасли 10.13, TestFirm).",
         12, False, MUTED)
    footer(s, 5)

    # 6 taxes
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "НАЛОГИ И ВЗНОСЫ", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.55), "45 млн ₽ в бюджет и фонды за 2025 год", 28, True, INK)
    taxes = [
        ("Налог на прибыль", "19,4 млн ₽"),
        ("НДС", "6,5 млн ₽"),
        ("Страховые взносы (СФР)", "19,0 млн ₽"),
        ("Транспортный налог", "0,04 млн ₽"),
        ("Итого уплачено", "44,9 млн ₽"),
    ]
    y = Inches(1.9)
    for i, (a, b) in enumerate(taxes):
        box(s, Inches(0.6), y, Inches(7.3), Inches(0.7), WHITE, RGBColor(0xE4, 0xDD, 0xD0))
        text(s, Inches(0.8), y + Inches(0.16), Inches(4.4), Inches(0.4), a, 16, i == 4, INK)
        text(s, Inches(5.3), y + Inches(0.16), Inches(2.4), Inches(0.4), b, 16, True, GREEN, PP_ALIGN.RIGHT)
        y += Inches(0.78)
    kpi(s, Inches(8.2), Inches(1.9), Inches(4.5), Inches(2.2), "×3,6",
        "рост налогов и взносов 2024→2025. ФНС: налоги +304%, взносы +218%. 2024 ≈ 12,4 млн ₽ (оценка по темпам).")
    kpi(s, Inches(8.2), Inches(4.3), Inches(4.5), Inches(2.0), "ОСНО",
        "Плательщик налога на прибыль и НДС. Просим инструмент под рост производства, не льготу «вместо» налогов.")
    footer(s, 6)

    # 7 base
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "ПРОИЗВОДСТВЕННАЯ БАЗА", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.55), "Площадка готова к следующему шагу мощности", 26, True, INK)
    items = [
        ("2 000 м²", "Производство: термообработка, копчение, сыровяление, формовка, нарезка, выпечка, шоковая заморозка."),
        ("1 300 м²", "Холодильный контур. Без него не держатся ни сети, ни экспорт, ни рост SKU."),
        ("300 м²", "Склад готовой продукции. Узкое место при росте отгрузок в сети и на АЗС."),
        ("110 человек", "Текущий штат. ССЧ ФНС: 9 → 32 → 64. Рост людей не догоняет рост тонн — нужна автоматизация упаковки."),
        ("68,9 млн ₽", "Основные средства на 31.12.2025. Лизинг уже используем. Следующий контур — роботы + Variovac."),
        ("Федресурс", "10 сообщений о лизинге. Опыт работы с лизинговыми инструментами уже есть."),
    ]
    for i, (n, l) in enumerate(items):
        col, row = i % 3, i // 3
        kpi(s, Inches(0.6 + col * 4.15), Inches(1.95 + row * 2.3), Inches(4.0), Inches(2.1), n, l)
    footer(s, 7)

    # 8 products
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "ПРОДУКЦИЯ И КАЧЕСТВО", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.55), "Три линейки, один халяльный контур", 28, True, INK)
    cols = [
        ("Заморозка", "Сосиски гриль, котлеты, пепперони и топпинги, мясные заготовки. Срок до 360 суток при −18 °C."),
        ("Охлаждёнка", "Сосиски и сардельки, колбасы, ветчины, казылык. Срок 30 суток. Фасовка в том числе по 8 шт."),
        ("Выпечка", "Эчпочмак, губадия, перемяч, самса, чак-чак, сосиска в тесте — АЗС и ритейл."),
    ]
    for i, (n, l) in enumerate(cols):
        kpi(s, Inches(0.6 + i * 4.15), Inches(1.95), Inches(4.0), Inches(2.6), n, l)
    box(s, Inches(0.6), Inches(4.8), Inches(12.1), Inches(1.55), WHITE, RGBColor(0xE4, 0xDD, 0xD0))
    text(s, Inches(0.85), Inches(5.0), Inches(11.6), Inches(1.2),
         ["Сертификаты только подтверждённые: Халяль ДУМ РТ № 614A/2024, HACCP, ISO 22000:2018, ТР ТС 021/2011.",
          "Прослеживаемость партий, ВСД «Меркурий», входной контроль сырья. Без ГМО и трансглютаминазы."],
         16, False, INK)
    footer(s, 8)

    # 9 market
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "РЫНОК", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.55), "Якорные каналы уже федеральные", 28, True, INK)
    left = [
        "АЗС «Татнефть» — сосиски и котлеты по сети.",
        "СМАРТЕН — халяль в прикассовой зоне АЗС.",
        "АО «ОМПК» — СТМ Aslam, традиционные колбасы.",
        "Ритейл: EuroSpar, Бэхетле, Metro, Мираторг.",
    ]
    right = [
        "HoReCa: GFC, «СвитЛайф»; корпоративное питание КВЗ.",
        "Отгрузка EXW Казань, регулярные плечи по РФ.",
        "Экспорт: RUB, USD, KZT, UZS, KGS, BYN, AZN.",
        "1 кв. 2026: 229 млн ₽ — темп к обороту 1 млрд.",
    ]
    y = Inches(1.95)
    for a, b in zip(left, right):
        box(s, Inches(0.6), y, Inches(6.0), Inches(0.85), WHITE, RGBColor(0xE4, 0xDD, 0xD0))
        text(s, Inches(0.8), y + Inches(0.22), Inches(5.6), Inches(0.5), a, 15)
        box(s, Inches(6.8), y, Inches(5.9), Inches(0.85), WHITE, RGBColor(0xE4, 0xDD, 0xD0))
        text(s, Inches(7.0), y + Inches(0.22), Inches(5.5), Inches(0.5), b, 15)
        y += Inches(0.98)
    text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.45),
         "2025: 869 млн ₽ — выход за порог малого предприятия (800 млн ₽, 209-ФЗ). Категория: среднее. Цель 2026 — свыше 1 млрд ₽.",
         13, False, MUTED)
    footer(s, 9)

    # 10 robots
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "ПРОЕКТ", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.55), "Роботизация линий: четыре контура", 28, True, INK)
    blocks = [
        ("01", "Палетоукладчики. Снять ручной штабель с охлаждёнки и заморозки. Тонны в смену, меньше травматизма."),
        ("02", "Раскладка сосисок по 8 шт. Автомат на машинах «Вариовак». Горлышко фасовки, не рецептуры."),
        ("03", "Машинное зрение. Геометрия, дефект оболочки, комплектность упаковки — вместо выборочного осмотра."),
        ("04", "Линия Variovac. Упаковочный контур, к которому крепятся робот раскладки и зрение. Без линии нет такта."),
    ]
    for i, (n, l) in enumerate(blocks):
        col, row = i % 2, i // 2
        kpi(s, Inches(0.6 + col * 6.3), Inches(1.95 + row * 2.15), Inches(6.1), Inches(2.0), n, l)
    footer(s, 10)

    # 11 barrier
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "БАРЬЕР", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.7), "Железо компенсируют. Интеграцию — почти нет", 26, True, INK)
    kpi(s, Inches(0.6), Inches(2.1), Inches(5.4), Inches(4.2), "×3",
        "Затраты на интеграцию роботов в действующее производство могут втрое превышать стоимость машин: инженерия, оснастка, синхронизация с Variovac, безопасность, простой, обучение.")
    rights = [
        "Льготный заём РФРП — до 200 млн ₽, 3–5%, до 5 лет.",
        "Нацпроект «Производительность труда»: ФРП от 3% + РЦК.",
        "Региональные субсидии на оборудование: 3 млрд ₽ в бюджете РТ на 2026.",
        "Льготный лизинг: до 40 млн ₽, 9,5% / 11,5%.",
    ]
    y = Inches(2.1)
    for r in rights:
        box(s, Inches(6.2), y, Inches(6.5), Inches(0.95), WHITE, RGBColor(0xE4, 0xDD, 0xD0))
        text(s, Inches(6.4), y + Inches(0.25), Inches(6.1), Inches(0.55), r, 15)
        y += Inches(1.05)
    footer(s, 11)

    # 12 ask
    s = blank(prs)
    add_logo(s)
    text(s, Inches(0.6), Inches(0.85), Inches(12), Inches(0.3), "ПРОСЬБА К МИНПРОМТОРГУ РТ", 13, True, GREEN)
    text(s, Inches(0.6), Inches(1.15), Inches(12), Inches(0.55), "Три поручения, которые двигают проект", 26, True, INK)
    asks = [
        "1. Подтвердить применимость мер к проекту — включая интеграцию, не только CAPEX на манипуляторы. Письмо Председателю Правительства РТ А. В. Песошину направлено.",
        "2. Организовать консультацию РЦК и Центра развития промышленной робототехники (Университет Иннополис) — дорожная карта на нашей площадке.",
        "3. Посадить за один стол Минпромторг, Минэкономики и Минсельхоз РТ: условия займа РФРП и субсидии 2026 года под конкретный перечень оборудования и работ.",
    ]
    y = Inches(1.95)
    for a in asks:
        box(s, Inches(0.6), y, Inches(12.1), Inches(1.25), WHITE, RGBColor(0xE4, 0xDD, 0xD0))
        text(s, Inches(0.85), y + Inches(0.25), Inches(11.6), Inches(0.85), a, 16)
        y += Inches(1.4)
    text(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.4),
         "Уже пользовались мерами поддержки МСП: 3,5 млн ₽, 23 меры, 0 нарушений. Готовы к контуру промышленной модернизации.",
         13, False, MUTED)
    footer(s, 12)

    # 13 close
    s = blank(prs, dark=True)
    add_logo(s, dark=True)
    text(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.35),
         "ГОТОВЫ К СЛЕДУЮЩЕЙ ВСТРЕЧЕ НА ПЛОЩАДКЕ", 13, True, GOLD)
    text(s, Inches(0.6), Inches(1.9), Inches(12), Inches(2.4),
         ["869 млн уже есть.", "1 млрд — следующий такт,", "если линия не останется ручной."],
         34, True, CREAM)
    kpi(s, Inches(0.6), Inches(5.15), Inches(4.0), Inches(1.25), "+7 987 217-02-02", "Адаева Гульнара Нурхаматовна", True)
    kpi(s, Inches(4.8), Inches(5.15), Inches(4.0), Inches(1.25), "info@kazandelikates.tatar", "kazandelikates.tatar", True)
    kpi(s, Inches(9.0), Inches(5.15), Inches(3.7), Inches(1.25), "Казань, Аграрная, 2", "Агропромпарк «Казань»", True)
    footer(s, 13, dark=True)

    prs.save(OUT)
    print(f"wrote {OUT} ({OUT.stat().st_size} bytes)")


if __name__ == "__main__":
    main()
